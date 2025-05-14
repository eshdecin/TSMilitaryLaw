import os
import fitz  # PyMuPDF for PDF
from docx import Document as DocxDocument
from pptx import Presentation
import pandas as pd

from pathlib import Path

def extract_text_from_pdf(file_path):
    text = ""
    doc = fitz.open(file_path)
    for page in doc:
        text += page.get_text()
    return text

def extract_text_from_docx(file_path):
    text = ""
    doc = DocxDocument(file_path)
    for para in doc.paragraphs:
        text += para.text + "\n"
    return text

def extract_text_from_txt(file_path):
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()

def extract_text_from_pptx(file_path):
    text = ""
    prs = Presentation(file_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

def extract_text_from_xlsx(file_path):
    text = ""
    try:
        xls = pd.ExcelFile(file_path)
        for sheet in xls.sheet_names:
            df = xls.parse(sheet)
            text += df.to_string(index=False) + "\n"
    except Exception as e:
        print(f"[!] Failed to read {file_path}: {e}")
    return text

def load_documents_from_folder(folder_path="pdfs"):
    all_docs = []

    for root, dirs, files in os.walk(folder_path):
        for filename in files:
            file_path = os.path.join(root, filename)
            ext = filename.lower()

            try:
                if ext.endswith(".pdf"):
                    print(f"[+] Loading PDF: {filename}")
                    text = extract_text_from_pdf(file_path)

                elif ext.endswith(".docx"):
                    print(f"[+] Loading DOCX: {filename}")
                    text = extract_text_from_docx(file_path)

                elif ext.endswith(".txt"):
                    print(f"[+] Loading TXT: {filename}")
                    text = extract_text_from_txt(file_path)

                elif ext.endswith(".pptx"):
                    print(f"[+] Loading PPTX: {filename}")
                    text = extract_text_from_pptx(file_path)

                elif ext.endswith(".xlsx"):
                    print(f"[+] Loading XLSX: {filename}")
                    text = extract_text_from_xlsx(file_path)

                else:
                    print(f"[x] Skipping unsupported file: {filename}")
                    continue

                all_docs.append({
                    "filename": filename,
                    "content": text
                })

            except Exception as e:
                print(f"[!] Error processing {filename}: {e}")

    return all_docs

if __name__ == "__main__":
    documents = load_documents_from_folder()
    print(f"\nTotal files loaded: {len(documents)}")
    for doc in documents:
        print(f"- {doc['filename']} ({len(doc['content'])} chars)")
